import streamlit as st
import json
import io
import os
import re
import cv2
import numpy as np

from PIL import Image
from paddleocr import PaddleOCR

import google.generativeai as genai

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# =========================================================
# CONFIG
# =========================================================

st.set_page_config(
    page_title="Sistema Master PRO",
    layout="wide"
)

st.title("🚀 Sistema Master PRO")
st.write("Imagem → PowerPoint Editável com IA")

# =========================================================
# API KEY
# =========================================================

GEMINI_KEY = st.secrets.get("GEMINI_API_KEY")

if not GEMINI_KEY:
    st.error("Configure GEMINI_API_KEY no secrets.toml")
    st.stop()

genai.configure(api_key=GEMINI_KEY)

# =========================================================
# OCR
# =========================================================

@st.cache_resource
def carregar_ocr():
    return PaddleOCR(
        use_angle_cls=True,
        lang='pt'
    )

ocr = carregar_ocr()

# =========================================================
# GEMINI MODEL
# =========================================================

def carregar_modelo():
    system_prompt = """
Você é um sistema profissional de reconstrução de layouts.

Sua função:
Transformar uma imagem em estrutura JSON perfeita para PowerPoint editável.

RETORNE APENAS JSON VÁLIDO.

NUNCA:
- explique
- use markdown
- use ```json

Detecte:
- textos
- imagens
- formas
- botões
- ícones
- blocos
- fundos

Regras:
- preserve alinhamento
- preserve proporções
- preserve hierarquia visual
- preserve todos os textos
- preserve cores
- preserve posição

Formato obrigatório:
[
  {
    "tipo":"texto",
    "conteudo":"EXEMPLO",
    "fonte":"Montserrat",
    "peso":"bold",
    "align":"center",
    "cor_hex":"#FFFFFF",
    "tamanho_fonte":28,
    "x_percent":0.1,
    "y_percent":0.2,
    "largura_percent":0.5,
    "altura_percent":0.1
  }
]

Tipos permitidos:
- texto
- forma
- imagem
"""

    return genai.GenerativeModel(
        model_name="gemini-2.5-flash",
        system_instruction=system_prompt
    )

# =========================================================
# UTILS
# =========================================================

def limpar_json(texto):
    texto = texto.replace("```json", "")
    texto = texto.replace("```", "").strip()

    match = re.search(r".*", texto, re.S)

    if not match:
        raise Exception("JSON inválido retornado pela IA.")

    return json.loads(match.group(0))

def hex_para_rgb(hex_color):
    try:
        hex_color = hex_color.replace("#", "")

        if len(hex_color) != 6:
            return RGBColor(0, 0, 0)

        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)

        return RGBColor(r, g, b)

    except:
        return RGBColor(0, 0, 0)

def detectar_textos_ocr(img_pil):
    img_cv = cv2.cvtColor(
        np.array(img_pil),
        cv2.COLOR_RGB2BGR
    )

    resultado = ocr.ocr(img_cv)

    textos = []

    h, w = img_cv.shape[:2]

    for bloco in resultado:
        for linha in bloco:
            pontos = linha[0]
            texto = linha[1][0]

            x1 = pontos[0][0]
            y1 = pontos[0][1]

            x2 = pontos[2][0]
            y2 = pontos[2][1]

            textos.append({
                "texto": texto,
                "x_percent": x1 / w,
                "y_percent": y1 / h,
                "largura_percent": (x2 - x1) / w,
                "altura_percent": (y2 - y1) / h
            })

    return textos

# =========================================================
# IA ANALYSIS
# =========================================================

def analisar_imagem(modelo, img_pil, textos_ocr):

    prompt = f"""
OCR DETECTADO:
{textos_ocr}

Agora gere o JSON final completo.
"""

    response = modelo.generate_content(
        [img_pil, prompt]
    )

    return limpar_json(response.text)

# =========================================================
# POWERPOINT
# =========================================================

def adicionar_texto(slide, elemento, larg_slide, alt_slide):

    left = Inches(elemento["x_percent"] * larg_slide)
    top = Inches(elemento["y_percent"] * alt_slide)

    width = Inches(elemento["largura_percent"] * larg_slide)
    height = Inches(elemento["altura_percent"] * alt_slide)

    textbox = slide.shapes.add_textbox(
        left,
        top,
        width,
        height
    )

    tf = textbox.text_frame
    p = tf.paragraphs[0]

    p.text = elemento.get("conteudo", "")

    # ALIGN
    align = elemento.get("align", "left")

    if align == "center":
        p.alignment = PP_ALIGN.CENTER

    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT

    else:
        p.alignment = PP_ALIGN.LEFT

    # FONT
    font = p.font

    font.name = elemento.get(
        "fonte",
        "Arial"
    )

    font.size = Pt(
        elemento.get(
            "tamanho_fonte",
            18
        )
    )

    font.color.rgb = hex_para_rgb(
        elemento.get(
            "cor_hex",
            "#000000"
        )
    )

    peso = elemento.get(
        "peso",
        ""
    ).lower()

    if peso == "bold":
        font.bold = True

def adicionar_forma(slide, elemento, larg_slide, alt_slide):

    left = Inches(elemento["x_percent"] * larg_slide)
    top = Inches(elemento["y_percent"] * alt_slide)

    width = Inches(elemento["largura_percent"] * larg_slide)
    height = Inches(elemento["altura_percent"] * alt_slide)

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height
    )

    shape.fill.solid()

    shape.fill.fore_color.rgb = hex_para_rgb(
        elemento.get(
            "cor_hex",
            "#CCCCCC"
        )
    )

    shape.line.color.rgb = hex_para_rgb(
        elemento.get(
            "cor_hex",
            "#CCCCCC"
        )
    )

def adicionar_imagem(
    slide,
    elemento,
    imagem,
    larg_slide,
    alt_slide
):

    img_w, img_h = imagem.size

    x = int(elemento["x_percent"] * img_w)
    y = int(elemento["y_percent"] * img_h)

    w = int(elemento["largura_percent"] * img_w)
    h = int(elemento["altura_percent"] * img_h)

    crop = imagem.crop((x, y, x + w, y + h))

    img_bytes = io.BytesIO()

    crop.save(
        img_bytes,
        format="PNG"
    )

    img_bytes.seek(0)

    slide.shapes.add_picture(
        img_bytes,
        Inches(elemento["x_percent"] * larg_slide),
        Inches(elemento["y_percent"] * alt_slide),
        Inches(elemento["largura_percent"] * larg_slide),
        Inches(elemento["altura_percent"] * alt_slide)
    )

def gerar_ppt(
    dados,
    imagem,
    widescreen=True
):

    prs = Presentation()

    if widescreen:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        larg_slide = 13.333
        alt_slide = 7.5

    else:
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(10)

        larg_slide = 10
        alt_slide = 10

    slide = prs.slides.add_slide(
        prs.slide_layouts[6]
    )

    for elemento in dados:

        tipo = elemento.get(
            "tipo",
            ""
        ).lower()

        try:

            if tipo == "texto":
                adicionar_texto(
                    slide,
                    elemento,
                    larg_slide,
                    alt_slide
                )

            elif tipo == "forma":
                adicionar_forma(
                    slide,
                    elemento,
                    larg_slide,
                    alt_slide
                )

            elif tipo == "imagem":
                adicionar_imagem(
                    slide,
                    elemento,
                    imagem,
                    larg_slide,
                    alt_slide
                )

        except Exception as e:
            print("Erro elemento:", e)

    output = io.BytesIO()

    prs.save(output)

    output.seek(0)

    return output

# =========================================================
# UI
# =========================================================

arquivo = st.file_uploader(
    "Envie uma imagem",
    type=["png", "jpg", "jpeg"]
)

if arquivo:

    img_pil = Image.open(arquivo).convert("RGB")

    st.image(
        img_pil,
        use_container_width=True
    )

    formato = st.selectbox(
        "Formato",
        [
            "Widescreen",
            "Quadrado"
        ]
    )

    if st.button(
        "🚀 Gerar PowerPoint",
        type="primary"
    ):

        try:

            with st.status(
                "Processando...",
                expanded=True
            ) as status:

                st.write("1️⃣ Extraindo OCR...")
                textos_ocr = detectar_textos_ocr(img_pil)

                st.write("2️⃣ Analisando layout IA...")
                modelo = carregar_modelo()

                dados = analisar_imagem(
                    modelo,
                    img_pil,
                    textos_ocr
                )

                st.write("3️⃣ Gerando PowerPoint...")
                ppt = gerar_ppt(
                    dados,
                    img_pil,
                    widescreen=(formato == "Widescreen")
                )

                status.update(
                    label="✅ Finalizado",
                    state="complete"
                )

            st.success("PowerPoint gerado!")

            st.download_button(
                label="⬇️ Baixar PPTX",
                data=ppt,
                file_name="master_pro.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

            with st.expander("JSON IA"):
                st.json(dados)

        except Exception as e:
            st.error(f"Erro: {e}")